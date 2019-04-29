from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Bar chart geratfffsgfggfagfagghhshhsgsgsssssssssssssssssssssssss"
# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest','north']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7, 45))

# add chart to slide --------------------
x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5.5)
chart = slide.shapes.add_chart(
	XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False
#chart.has_text_frame = True
#chart.text_frame = "Bar graph"
chart.has_title = True
chart.chart_title.has_text_frame = True
chart.chart_title.format.fill.background()
chart.chart_title.text_frame.auto_size
parag = chart.chart_title.text_frame.add_paragraph()
parag.text = "Bar tiltele \n Graph"
prs.save('chart-01.pptx')
