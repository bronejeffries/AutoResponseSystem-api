from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(5.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tfp = tf.add_paragraph()
tfp.text = "This is text inside a textbox"
tfp.level = 1
tfp.font.size = Pt(40)

for n in range(12):
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's are we going to embarrass you bold"
    p.level = 2
    p.font.size = Pt(20)

blank_slide_layout2 = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_slide_layout2)
txBox2 = slide2.shapes.add_textbox(left, top, width, height)
tf2 = txBox2.text_frame
tf2.text = "This is text inside a textbox"
p2 = tf2.add_paragraph()
p2.text = "This is a second paragraph that's bold"
p2.font.bold = True
p2 = tf2.add_paragraph()
p2.text = "This is a third paragraph that's big"
p2.font.size = Pt(40)
prs.save('test.pptx')
