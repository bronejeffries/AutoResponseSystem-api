import os
from django.http import HttpResponse
from django.shortcuts import render
from pptx import Presentation
from Ars.decorators import Get_check,Post_check
import io, random
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from Ars.models import Session, Topic, Question
# from .categoryViews import get_Question

def generate_random():
    return random.randint(1, 1001)

def get_Question(pk):
    question = None
    try:
        question = Question.objects.get(id=pk)
    except Question.DoesNotExist:
        return False
    else:
        return question

def question_options(question):
     options = question.option_set.all()
     return options

def get_Topic(id):
    topic = None
    try:
        topic = Topic.objects.get(id=id)
    except Topic.DoesNotExist:
        return False
    else:
        return topic

def topic_comments(topic):
    comments = topic.comment_set.all()
    return comments

def makeQuestionPresentation(question_id, name = None):
    # path
    question = get_Question(question_id)
    presentation_name = ""
    if name is not None:
        presentation_name = name
    else:
        presentation_name = str(question_id)+ "que_" + str(generate_random())

    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    storage_path = os.path.join(BASE_DIR, 'moderator/presentations')
    presentation_path = os.path.join(storage_path, presentation_name + '.pptx')
    # create presentation with one slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"({question.question})\t\n"

    # define chart bar data
    chart_data = CategoryChartData()

    question_options_choices = None
    question_options_list = None
    if question:
        question_options_list = question_options(question)
        question_options_choices = [int(option.choices) for option in question_options_list]
    max_length_cat = 10
    chart_data.categories = [ str(option.option[0:max_length_cat]) for option in question_options_list ]
    chart_data.add_series('series 1', tuple(question_options_choices))

    # add chart to slide
    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    prs.save(presentation_path)
    return (prs , presentation_name)


def makeTopicpresentation(topic_id, name = None):
    topic = get_Topic(topic_id)
    presentation_name = ""
    if name is not None:
        presentation_name = name
    else:
        presentation_name = f"{topic_id}top_{generate_random()}"

    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    storage_path = os.path.join(BASE_DIR, 'moderator/presentations')
    presentation_path = os.path.join(storage_path, presentation_name + '.pptx')

    prs = Presentation()
    left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(5.5)
    blank_slide_layout = prs.slide_layouts[6]

    # create first slide with title
    slide = prs.slides.add_slide(blank_slide_layout)

    # add textbox to first slide
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    # title paragraph
    tfp = tf.add_paragraph()
    tfp.text = str(topic.topic_name)
    tfp.level = 1
    tfp.font.size = Pt(40)

    SLIDE_MAX_CONTENT = 12
    COUNT = 0;
    topic_comment_s = topic_comments(topic)

    for comment in topic_comment_s:
        if COUNT < SLIDE_MAX_CONTENT :
            # populate the current slide

            p = tf.add_paragraph()
            p.text = str(comment.comment)
            p.level = 2
            p.font.size = Pt(20)

        else:
            # add new slide
            COUNT = 0
            slide = prs.slides.add_slide(blank_slide_layout)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = str(comment.comment)
            p.level = 2
            p.font.size = Pt(20)

        COUNT +=1

    prs.save(presentation_path)
    return (prs, presentation_name)

@Get_check
def getPresentation(request,name):
    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    storage_path = os.path.join(BASE_DIR, 'moderator/presentations')
    presentation_path = os.path.join(storage_path, name + '.pptx')
    # file = open(presentation_path)
    # prs = Presentation(presentation_path)
    # file.close()
    targetstream = io.BytesIO()
    with open(presentation_path,'rb') as f:
        targetstream = io.BytesIO(f.read())
    # prs.save(targetstream)
    response = HttpResponse(content_type='application/vnd.ms-powerpoint')
    response['Content-Disposition'] = 'attachment; filename="'+name+'.pptx"'
    response.write(targetstream.getvalue())
    targetstream.close()
    return response

@Get_check
def viewPresentation(request,name):
    scheme = request.is_secure() and "https" or "http"
    pathheader = scheme+"://"+request.META['HTTP_HOST']
    return render(request,'moderator/viewpresentation.html',{'presentationname':name,'pathheader':pathheader})
